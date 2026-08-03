import io

import pytest
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from app.bell_tasks.generator import NUM_BOXES, NUM_SLIDES, QUESTIONS_PER_TOPIC, generate_bell_tasks_pptx
from app.bell_tasks.layout import BOX_TO_ROW_COL, box_bounds
from app.core.errors import TopicNotFoundError

# A deliberate mix: some diagram-bearing topics, some not.
SIX_TOPIC_IDS = [
    "angles_triangle",  # diagram
    "area_rectangle",  # diagram
    "fractions_add_subtract",  # no diagram, always produces a fraction
    "linear_two_step",  # no diagram
    "probability_single_event",  # diagram
    "bar_chart_construct",  # diagram, sometimes long prompt
]

GRID_TABLE_SHAPE_NAME = "Google Shape;146;p3"

_M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"


def _grid_table(slide):
    for shape in slide.shapes:
        if shape.has_table and shape.name == GRID_TABLE_SHAPE_NAME:
            return shape.table
    raise AssertionError("grid table not found on slide")


def _cell_full_text(cell) -> str:
    """Flattens a cell's paragraph back to text, including native-equation
    content (read from each embedded equation's own mc:Fallback plain text,
    since that's the only place the OMML content also exists as plain text) -
    a plain `run.text` join would silently skip any fraction/exponent, since
    those are no longer plain `<a:r>` runs at all."""
    parts = []
    for child in cell.text_frame.paragraphs[0]._p:
        tag = etree.QName(child).localname
        if tag == "r":
            texts = child.findall(f".//{{{_A_NS}}}t")
            parts.append("".join(t.text or "" for t in texts))
        elif tag == "AlternateContent":
            fallback_texts = child.findall(f".//{{{_MC_NS}}}Fallback//{{{_A_NS}}}t")
            parts.append("".join(t.text or "" for t in fallback_texts))
    return "".join(parts)


def test_generates_five_slides_with_six_distinct_boxes_each():
    data = generate_bell_tasks_pptx(SIX_TOPIC_IDS)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == NUM_SLIDES

    for slide in prs.slides:
        table = _grid_table(slide)
        for box in range(1, NUM_BOXES + 1):
            row, col = BOX_TO_ROW_COL[box]
            cell = table.cell(row, col)
            text = _cell_full_text(cell)
            assert text.startswith(f"{box}. ")
            assert len(text) > len(f"{box}. ")


def _box_picture_blobs(slide, box) -> tuple:
    """Raw image bytes of every picture shape positioned inside `box`'s own
    cell bounds on this slide - `diagram_rect` always places a box's diagram
    picture within `box_bounds(box)`, so geometric containment reliably
    matches a picture to its box without needing to intercept generation
    order."""
    bx, by, bw, bh = box_bounds(box)
    blobs = []
    for shape in slide.shapes:
        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            continue
        if bx <= shape.left < bx + bw and by <= shape.top < by + bh:
            blobs.append(shape.image.blob)
    return tuple(blobs)


def test_each_box_shows_five_distinct_questions_across_the_week():
    data = generate_bell_tasks_pptx(SIX_TOPIC_IDS)
    prs = Presentation(io.BytesIO(data))

    for box in range(1, NUM_BOXES + 1):
        row, col = BOX_TO_ROW_COL[box]
        seen = set()
        for slide in prs.slides:
            table = _grid_table(slide)
            cell = table.cell(row, col)
            text = _cell_full_text(cell)
            # Text alone doesn't always prove two questions differ: some
            # topics move their distinguishing numbers onto the diagram and
            # leave the prompt itself generic (e.g. area_rectangle, after
            # this project's own prose-stripping pass), so two genuinely
            # different questions can share identical prompt text. Fold in
            # each box's own diagram image bytes too, so distinctness is
            # judged on what's actually rendered, not just the text run.
            key = (text, _box_picture_blobs(slide, box))
            seen.add(key)
        assert len(seen) == QUESTIONS_PER_TOPIC, f"box {box} did not get {QUESTIONS_PER_TOPIC} distinct questions"


def test_runs_use_calibri_for_words_and_cambria_math_for_numbers():
    data = generate_bell_tasks_pptx(SIX_TOPIC_IDS)
    prs = Presentation(io.BytesIO(data))

    seen_fonts = set()
    for slide in prs.slides:
        table = _grid_table(slide)
        for box in range(1, NUM_BOXES + 1):
            row, col = BOX_TO_ROW_COL[box]
            cell = table.cell(row, col)
            for run in cell.text_frame.paragraphs[0].runs:
                assert run.font.size is not None
                assert run.font.size.pt == 18
                seen_fonts.add(run.font.name)

    assert "Calibri" in seen_fonts
    assert "Cambria Math" in seen_fonts


def test_fraction_topic_produces_a_real_native_equation_object():
    # fractions_add_subtract always produces a standalone fraction in its
    # prompt - confirms the real <m:f> equation object actually gets emitted
    # into the deck, not just plain Cambria Math text.
    data = generate_bell_tasks_pptx(SIX_TOPIC_IDS)
    prs = Presentation(io.BytesIO(data))

    box_for_fractions = SIX_TOPIC_IDS.index("fractions_add_subtract") + 1
    row, col = BOX_TO_ROW_COL[box_for_fractions]
    found_fraction = False
    for slide in prs.slides:
        table = _grid_table(slide)
        cell = table.cell(row, col)
        p_xml = cell.text_frame.paragraphs[0]._p
        if p_xml.findall(f".//{{{_M_NS}}}f"):
            found_fraction = True
            break
    assert found_fraction, "expected at least one real <m:f> equation across the week"


def test_pictures_stay_within_their_own_cell_bounds():
    data = generate_bell_tasks_pptx(SIX_TOPIC_IDS)
    prs = Presentation(io.BytesIO(data))

    for slide in prs.slides:
        pictures = [s for s in slide.shapes if s.shape_type == 13]
        # Every slide always has the 2 static template logos; anything beyond
        # that is a diagram this run added, and must sit inside *some* cell.
        added_pictures = pictures[2:] if len(pictures) > 2 else []
        for picture in added_pictures:
            matched = False
            for box in range(1, NUM_BOXES + 1):
                cell_left, cell_top, cell_width, cell_height = box_bounds(box)
                if (
                    picture.left >= cell_left
                    and picture.top >= cell_top
                    and picture.left + picture.width <= cell_left + cell_width
                    and picture.top + picture.height <= cell_top + cell_height
                ):
                    matched = True
                    break
            assert matched, f"picture at ({picture.left},{picture.top}) doesn't fit inside any cell"


def test_rejects_wrong_number_of_topics():
    with pytest.raises(ValueError):
        generate_bell_tasks_pptx(SIX_TOPIC_IDS[:5])


def test_rejects_duplicate_topics():
    duplicated = SIX_TOPIC_IDS[:5] + [SIX_TOPIC_IDS[0]]
    with pytest.raises(ValueError):
        generate_bell_tasks_pptx(duplicated)


def test_rejects_unknown_topic_id():
    bad = SIX_TOPIC_IDS[:5] + ["not_a_real_topic"]
    with pytest.raises(TopicNotFoundError):
        generate_bell_tasks_pptx(bad)
