import io

from lxml import etree
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

from app.bell_tasks import omml
from app.bell_tasks.generator import _template_path

_MC_NS = "http://schemas.openxmlformats.org/markup-compatibility/2006"
_A14_NS = "http://schemas.microsoft.com/office/drawing/2010/main"
_M_NS = "http://schemas.openxmlformats.org/officeDocument/2006/math"
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"

PURPLE = RGBColor(0x53, 0x1D, 0x60)


def _new_paragraph():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    return prs, textbox.text_frame.paragraphs[0]


def _roundtrip(prs) -> Presentation:
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return Presentation(buf)


def _find(el, path):
    return el.findall(path, namespaces={"mc": _MC_NS, "a14": _A14_NS, "m": _M_NS, "a": _A_NS})


def test_add_fraction_builds_a_real_oMath_fraction_element():
    prs, paragraph = _new_paragraph()
    omml.add_fraction(paragraph, "", "3", "4", 18.0, PURPLE)

    p_xml = paragraph._p
    fractions = _find(p_xml, ".//m:f")
    assert len(fractions) == 1
    numerators = _find(fractions[0], "./m:num/m:r/m:t")
    denominators = _find(fractions[0], "./m:den/m:r/m:t")
    assert numerators[0].text == "3"
    assert denominators[0].text == "4"


def test_add_fraction_applies_sign_to_the_numerator():
    prs, paragraph = _new_paragraph()
    omml.add_fraction(paragraph, "-", "3", "4", 18.0, PURPLE)

    numerators = _find(paragraph._p, ".//m:f/m:num/m:r/m:t")
    assert numerators[0].text == "-3"


def test_add_fraction_includes_a_plain_text_fallback():
    prs, paragraph = _new_paragraph()
    omml.add_fraction(paragraph, "", "3", "4", 18.0, PURPLE)

    fallback_text = _find(paragraph._p, ".//mc:Fallback//a:t")
    assert fallback_text[0].text == "3/4"


def test_add_exponent_builds_a_real_oMath_superscript_element():
    prs, paragraph = _new_paragraph()
    omml.add_exponent(paragraph, "x", "2", 18.0, PURPLE)

    p_xml = paragraph._p
    superscripts = _find(p_xml, ".//m:sSup")
    assert len(superscripts) == 1
    base = _find(superscripts[0], "./m:e/m:r/m:t")
    exponent = _find(superscripts[0], "./m:sup/m:r/m:t")
    assert base[0].text == "x"
    assert exponent[0].text == "2"


def test_add_fractional_exponent_nests_a_fraction_inside_the_superscript():
    prs, paragraph = _new_paragraph()
    omml.add_fractional_exponent(paragraph, "x", "1", "4", 18.0, PURPLE)

    p_xml = paragraph._p
    superscripts = _find(p_xml, ".//m:sSup")
    assert len(superscripts) == 1
    base = _find(superscripts[0], "./m:e/m:r/m:t")
    assert base[0].text == "x"
    nested_fraction_num = _find(superscripts[0], "./m:sup/m:f/m:num/m:r/m:t")
    nested_fraction_den = _find(superscripts[0], "./m:sup/m:f/m:den/m:r/m:t")
    assert nested_fraction_num[0].text == "1"
    assert nested_fraction_den[0].text == "4"


def test_math_run_formatting_carries_size_and_colour():
    prs, paragraph = _new_paragraph()
    omml.add_fraction(paragraph, "", "3", "4", 18.0, PURPLE)

    a_rprs = _find(paragraph._p, ".//m:r/m:rPr/a:rPr")
    assert len(a_rprs) == 2  # one for numerator run, one for denominator run
    for a_rpr in a_rprs:
        assert a_rpr.get("sz") == "1800"
        srgb = _find(a_rpr, "./a:solidFill/a:srgbClr")
        assert srgb[0].get("val") == "531D60"
        latin = _find(a_rpr, "./a:latin")
        assert latin[0].get("typeface") == "Cambria Math"


def test_equation_survives_a_real_save_and_reload_round_trip():
    prs, paragraph = _new_paragraph()
    omml.add_fraction(paragraph, "", "3", "4", 18.0, PURPLE)
    omml.add_exponent(paragraph, "x", "2", 18.0, PURPLE)

    reopened = _roundtrip(prs)
    reopened_p = reopened.slides[0].shapes[0].text_frame.paragraphs[0]._p
    assert len(_find(reopened_p, ".//m:f")) == 1
    assert len(_find(reopened_p, ".//m:sSup")) == 1


def test_equation_is_inserted_before_a_pre_existing_end_para_rpr():
    # Regression test for a real bug: found via visual inspection of an
    # actual rendered slide, not caught by any test written in advance, since
    # the XML is well-formed either way - PowerPoint just silently drops
    # anything positioned after `<a:endParaRPr>` rather than erroring. A
    # freshly cleared real template cell (exactly what generator.py's
    # `_set_cell_content` operates on) leaves a trailing `endParaRPr` behind,
    # which a naive `etree.SubElement(paragraph_xml, ...)` append would land
    # after - this asserts the fix (inserting before it) actually holds
    # against a real template cell, not just a synthetic paragraph.
    prs = Presentation(_template_path())
    slide = prs.slides[0]
    grid = next(s for s in slide.shapes if s.has_table and s.name == "Google Shape;146;p3")
    cell = grid.table.cell(0, 0)
    cell.text_frame.clear()
    paragraph = cell.text_frame.paragraphs[0]

    p_xml = paragraph._p
    assert p_xml.find(f".//{{{_A_NS}}}endParaRPr") is not None, (
        "test precondition failed: a freshly cleared template cell no longer "
        "leaves a trailing endParaRPr - the bug this test guards may no longer "
        "be reachable this way, but re-check rather than assume it's gone"
    )

    run = paragraph.add_run()
    run.text = "Solve "
    omml.add_exponent(paragraph, "x", "2", 18.0, PURPLE)

    tags = [etree.QName(c).localname for c in p_xml]
    assert tags.index("AlternateContent") < tags.index("endParaRPr")
    assert tags[-1] == "endParaRPr"


def test_multiple_equations_and_plain_runs_coexist_in_one_paragraph():
    prs, paragraph = _new_paragraph()
    run = paragraph.add_run()
    run.text = "Simplify "
    omml.add_fraction(paragraph, "", "3", "4", 18.0, PURPLE)
    run2 = paragraph.add_run()
    run2.text = " fully."

    p_xml = paragraph._p
    # Children should be: a:r, mc:AlternateContent, a:r, in that document order.
    children = list(p_xml)
    tags = [etree.QName(c).localname for c in children]
    assert tags.count("r") == 2
    assert tags.count("AlternateContent") == 1
    assert tags.index("AlternateContent") == 1
