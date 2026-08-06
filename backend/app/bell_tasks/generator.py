"""Generates a Bell Tasks .pptx: 6 chosen topics, one permanently assigned to
each of the reference template's 6 numbered boxes, with that topic's 5
generated questions (one per weekday slide) filling that box across all 5
slides. Reuses the existing verified generation pipeline end to end
(`build_worksheet`) - this module only adds a new *renderer* (pptx instead of
PDF) on top of it, no new question-generation logic.
"""

import importlib.resources
import io
import random

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

from app.bell_tasks import omml
from app.bell_tasks.diagram_raster import rasterize_drawing
from app.bell_tasks.layout import BOX_TO_ROW_COL, box_bounds, diagram_rect
from app.bell_tasks.math_tokenizer import ExponentSpan, FractionalExponentSpan, FractionSpan, TextSpan, tokenize
from app.core.registry import get_topic
from app.pdf.diagrams import render_diagram
from app.worksheet.builder import build_worksheet

NUM_SLIDES = 5
QUESTIONS_PER_TOPIC = 5
NUM_BOXES = 6
PURPLE = RGBColor(0x53, 0x1D, 0x60)
FONT_SIZE = Pt(18)
FONT_SIZE_PT = 18.0

_GRID_TABLE_SHAPE_NAME = "Google Shape;146;p3"


def _template_path() -> str:
    return str(importlib.resources.files("app.bell_tasks") / "assets" / "bell_task_template.pptx")


def _grid_table(slide):
    for shape in slide.shapes:
        if shape.has_table and shape.name == _GRID_TABLE_SHAPE_NAME:
            return shape.table
    raise RuntimeError("Bell Tasks template slide is missing its 6-box question grid table")


def _set_cell_content(cell, box: int, prompt: str) -> None:
    text_frame = cell.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]

    tokens: list = [TextSpan(f"{box}. ", "Calibri"), *tokenize(prompt)]
    for token in tokens:
        if isinstance(token, TextSpan):
            run = paragraph.add_run()
            run.text = token.text
            run.font.name = token.font
            run.font.size = FONT_SIZE
            run.font.color.rgb = PURPLE
        elif isinstance(token, FractionSpan):
            omml.add_fraction(paragraph, token.sign, token.numerator, token.denominator, FONT_SIZE_PT, PURPLE)
        elif isinstance(token, ExponentSpan):
            omml.add_exponent(paragraph, token.base, token.exponent, FONT_SIZE_PT, PURPLE)
        elif isinstance(token, FractionalExponentSpan):
            omml.add_fractional_exponent(
                paragraph, token.base, token.numerator, token.denominator, FONT_SIZE_PT, PURPLE
            )


def generate_bell_tasks_pptx(topic_ids: list[str]) -> bytes:
    if len(topic_ids) != NUM_BOXES:
        raise ValueError(f"Bell Tasks needs exactly {NUM_BOXES} topic ids, got {len(topic_ids)}")
    if len(set(topic_ids)) != NUM_BOXES:
        raise ValueError("Bell Tasks topic ids must all be distinct")

    rng = random.Random()

    questions_by_box: dict[int, tuple] = {}
    for box, topic_id in enumerate(topic_ids, start=1):
        topic = get_topic(topic_id)
        worksheet = build_worksheet(topic_id, topic.fixed_tier, count=QUESTIONS_PER_TOPIC, rng=rng)
        questions_by_box[box] = worksheet.questions

    prs = Presentation(_template_path())

    for slide_idx in range(NUM_SLIDES):
        slide = prs.slides[slide_idx]
        table = _grid_table(slide)

        for box in range(1, NUM_BOXES + 1):
            row, col = BOX_TO_ROW_COL[box]
            question = questions_by_box[box][slide_idx]

            cell = table.cell(row, col)
            _set_cell_content(cell, box, question.prompt)

            if question.diagram is not None:
                drawing = render_diagram(question.diagram)
                rect = diagram_rect(box_bounds(box), question.prompt, drawing.width, drawing.height)
                if rect is not None:
                    png_bytes = rasterize_drawing(drawing)
                    left, top, width, height = rect
                    slide.shapes.add_picture(io.BytesIO(png_bytes), left, top, width, height)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
